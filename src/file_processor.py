######################normal code#######################


# import os
# from pathlib import Path
# from pypdf import PdfReader
# from docx import Document as DocxDocument
# import pptx

# def extract_text_from_pdf(file_path: str) -> str:
#     """Extract text from a PDF file."""
#     try:
#         reader = PdfReader(file_path)
#         text = ""
#         for page in reader.pages:
#             page_text = page.extract_text()
#             if page_text:
#                 text += page_text + "\n"
#         return text
#     except Exception as e:
#         print(f"Error reading PDF {file_path}: {e}")
#         return ""

# def extract_text_from_docx(file_path: str) -> str:
#     """Extract text from a DOCX file."""
#     try:
#         doc = DocxDocument(file_path)
#         return "\n".join(para.text for para in doc.paragraphs)
#     except Exception as e:
#         print(f"Error reading DOCX {file_path}: {e}")
#         return ""

# def extract_text_from_pptx(file_path: str) -> str:
#     """Extract text from a PPTX file."""
#     try:
#         prs = pptx.Presentation(file_path)
#         text = ""
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#         return text
#     except Exception as e:
#         print(f"Error reading PPTX {file_path}: {e}")
#         return ""

# def extract_text_from_txt(file_path: str) -> str:
#     """Extract text from a plain text file."""
#     try:
#         with open(file_path, "r", encoding="utf-8") as f:
#             return f.read()
#     except Exception as e:
#         print(f"Error reading TXT {file_path}: {e}")
#         return ""

# def extract_text_from_file(file_path: str) -> str:
#     """Detect file extension and extract text accordingly."""
#     ext = Path(file_path).suffix.lower()
#     if ext == ".pdf":
#         return extract_text_from_pdf(file_path)
#     elif ext == ".docx":
#         return extract_text_from_docx(file_path)
#     elif ext == ".pptx":
#         return extract_text_from_pptx(file_path)
#     elif ext in [".txt", ".md", ".csv", ".json"]:
#         return extract_text_from_txt(file_path)
#     else:
#         print(f"Unsupported file type for {file_path}")
#         return ""


#################multi-modal code###################

import os
from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
import pptx
from PIL import Image
import pytesseract

def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = DocxDocument(file_path)
        return "\n".join(para.text for para in doc.paragraphs)
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return ""

def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"Error reading TXT {file_path}: {e}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error reading image {file_path}: {e}")
        return ""

def extract_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext in [".txt", ".md", ".csv", ".json"]:
        return extract_text_from_txt(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
        return extract_text_from_image(file_path)
    else:
        print(f"Unsupported file type for {file_path}")
        return ""


